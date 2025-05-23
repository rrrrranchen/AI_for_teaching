import uuid
from openai import OpenAI
import time
import os


# 定义自己的 API Key
key = 'sk-b7550aa67ed840ffacb5ca051733802c'
api_url = "https://api.deepseek.com"  # DeepSeek 的 API 地址
# OpenAI 参数设置：API Key + API Interface (这里访问接口为 DeepSeek 的 API 地址)


# 逐字打印效果
def printChar(text, delay=0.1):
    for char in text:
        print(char, end='', flush=True)  # 使用 end='' 防止自动换行，flush=True 确保立即打印
        time.sleep(delay)
    print()  # 最后打印一个换行符


# 发送请求到 DeepSeek
def sendToDeepSeek(say):
    print('正在验证身份，请稍等....')
    # 请求接口并验证身份，创建客户端对象
    client = OpenAI(api_key=key, base_url=api_url)
    print('正在思考，请耐心等待...')
    # 发送请求数据并等待获取响应数据
    response = client.chat.completions.create(
        model="deepseek-chat",  # 使用的模型
        messages=[
            {"role": "system", "content": "你是一个专业的客服助手，请用正式的语气回答用户的问题。"},
            # {"role": "system", "content": "你是风趣幽默的客服，请用轻松幽默的语气回答用户的问题。"},
            {"role": "user", "content": say},
        ],
        stream=False  # 是否启用流式输出
    )
    # print(response)  # 如果需要调试，可以打印完整的响应
    return response.choices[0].message.content  # 返回模型生成的回复内容


# DeepSeek 问答环节
def callDeepSeek():
    # 主循环
    while True:
        myin = input('您请说：')  # 获取用户输入
        if myin == 'bye':  # 如果用户输入 "bye"，退出循环
            print('欢迎下次使用！再见！')
            break
        resp = sendToDeepSeek(myin)  # 发送用户输入到 DeepSeek 并获取回复
        printChar(resp)  # 逐字打印回复内容
        # print(resp)  # 如果需要直接打印完整回复，可以使用这行代码
        print('-----------------------------------------------------------')


# 测试接口
# callDeepSeek()


from pptx import Presentation, util
from datetime import datetime
from openai import OpenAI
from pptx import Presentation, util
from pptx.dml.color import RGBColor
import ast


def set_text_font(shape, font_name=None, font_size=None, font_color=None, font_bold=None, font_italic=None):
    text_frame = shape.text_frame

    # 遍历所有段落和文本运行
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            # 设置字体名称
            if font_name is not None:
                run.font.name = font_name
            # 设置字体大小（单位：磅）
            if font_size is not None:
                run.font.size = util.Pt(font_size)
            # 设置颜色
            if font_color is not None:  # NoneColor Type
                run.font.color.rgb = font_color
            # 设置加黑
            if font_bold is not None:  # False强制取消粗体, None继承父版样式
                run.font.bold = font_bold
            # 设置斜体
            if font_italic is not None:
                run.font.italic = font_italic
            # -> False显式取消加粗、斜体 -> if False... False & None 区别


def keep_text_font(shape, new_text):
    text_frame = shape.text_frame

    for paragraph in text_frame.paragraphs:
        # 创建新 Run 并继承格式（取第一个 Run 的格式）
        if paragraph.runs:
            # 清空所有 Run 的文本
            for run in paragraph.runs:
                run.text = ""

            first_run = paragraph.runs[0]
            new_run = paragraph.add_run()
            new_run.text = new_text

            # 复制字体属性
            new_run.font.name = first_run.font.name
            new_run.font.size = first_run.font.size
            new_run.font.bold = first_run.font.bold
            new_run.font.italic = first_run.font.italic
            # new_run.font.color.rgb = first_run.font.color.rgb
            # 检查颜色是否有效
            # if first_run.font.color is not None:  # and first_run.font.color.rgb -> 存在NoneType No RGB
            if (
                first_run.font.color is not None
                and hasattr(first_run.font.color, 'rgb')  # 关键检查
            ):
                new_run.font.color.rgb = first_run.font.color.rgb
            else:
                # 若颜色未设置，使用默认黑色或跳过
                new_run.font.color.rgb = RGBColor(0, 0, 0)  # 显式设置为黑色


def set_ppt_content(content, template, name, time, title, subtitle=None, keep_font=False,
                    font_name=None, font_size=None, font_bold=None, font_italic=None, font_color=None):
    prs = Presentation(template)

    # PPT模板的shapes全为文本框，这里可以省去判断步骤
    for slide in prs.slides:
        for shape in slide.shapes:
            if not (shape.has_text_frame and shape.text):
                continue
            text = shape.text
            # 设置标题 -> 防止副标题、小标题字样干扰
            if text == '标题':
                keep_text_font(shape, title)
                # if not keep_font: -> 设置字体样式... 这里保持为原格式
                set_text_font(shape, font_name='标准粗黑')
            # 设置副标题
            if text == '副标题':
                keep_text_font(shape, subtitle)
                set_text_font(shape, font_name='标准粗黑')
            # 设置老师姓名
            if '姓名' in text:
                keep_text_font(shape, f'主讲人：{name}')
                set_text_font(shape, font_name='等线')
            # 设置汇报日期
            if '日期' in text or '年月日' in text:
                keep_text_font(shape, f'时间：{time}')
                set_text_font(shape, font_name='等线')
            # 设置内容页
            if '小标题' in text:
                # shape.text = get_content(shape, content)
                # keep_text_font(shape, get_content(shape, content))
                sub_content, if_content = get_content(shape, content)
                keep_text_font(shape, sub_content)
                if not if_content:
                    set_text_font(shape, font_name='等线', font_bold=True)
            if not keep_font:
                set_text_font(shape, font_name, font_size, font_color, font_bold, font_italic)
    return prs


def get_content(shape, content):
    text = shape.text
    text = text.replace('小标题', '')

    if '概述' in text:
        i = int(text.replace('概述', '')) - 1
        return content[list(content.keys())[i]]['概述'], None  # False

    if_content = False

    if '内容' in text:
        text = text.replace('内容', '')
        if_content = True

    num = int(text.replace('.', ''))
    i, j, k = int(num / 100) - 1, int(num / 10) % 10 - 1, num % 10 - 1
    if j < 0:
        # return content.keys[k]
        key_lst = [k for k in content.keys()]
        return key_lst[k], if_content
    elif i < 0:
        # return content[content.keys[j]].keys[k]
        key_lst1 = list(content.keys())
        key_lst2 = [k for k in content[key_lst1[j]].keys() if k != '概述']  # list(content[key_lst1[j]].keys())
        return key_lst2[k], if_content
    # else:
    key_lst1 = list(content.keys())
    key_lst2 = [k for k in content[key_lst1[i]].keys() if k != '概述']
    key_lst3 = list(content[key_lst1[i]][key_lst2[j]].keys())
    # print(num, i+1, j+1, k+1, key_lst3)
    if if_content:
        return content[key_lst1[i]][key_lst2[j]][key_lst3[k]], if_content
    return key_lst3[k], if_content  # 输出if_content->调节小标题字体
    #     return content[content[content.keys[i]].keys[j]].values[k]
    # return content[content[content.keys[i]].keys[j]].keys[k]


project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPTTEMPLATE_FOLDER = os.path.join(project_root, 'static', 'template')
PPT_FOLDER=os.path.join(project_root,'static','uploads','ppts')
filename='template.pptx'
templatepath=os.path.join(PPTTEMPLATE_FOLDER, filename)
print(templatepath)
structure = {
'template.pptx': [
    [3, 3],
    [3, 3, 3],
    [3, 4, 3],
    [3, 3, 3, 3],
    [3, 3, 3]
    ]
}
def parse_teaching_plan(teaching_plan):
    """
    通过 AI 接口解析教学计划，提取学科（subject）和章节总知识点（chapter）。
    
    Args:
        teaching_plan (str): 教学计划的文本内容。
    
    Returns:
        tuple: (subject, chapter)，其中 subject 是学科，chapter 是章节总知识点。
    """
    print("正在通过 AI 接口解析教学计划...")
    
    # 更明确的提示词，要求AI返回特定格式
    prompt = f"""请从以下教学计划中提取信息并以严格的JSON格式返回：
    - "subject": 课程主题或学科名称
    - "chapter": 主要知识点（一个总的知识点）
    
    教学计划内容：
    {teaching_plan}
    
    返回格式示例：
    {{
        "subject": "计算机网络",
        "chapter": "网络分类、拓扑结构、TCP/IP协议"
    }}
    """
    
    try:
        # 调用 AI 接口解析教学计划
        ai_response = sendToDeepSeek(prompt)
        
        # 预处理响应：替换中文标点为英文标点
        ai_response = ai_response.replace('：', ':').replace('，', ',').replace('“', '"').replace('”', '"')
        
        # 尝试提取JSON部分（有时AI会在JSON前后添加解释文本）
        start_idx = ai_response.find('{')
        end_idx = ai_response.rfind('}') + 1
        json_str = ai_response[start_idx:end_idx]
        
        # 解析JSON
        ai_result = ast.literal_eval(json_str)
        subject = ai_result.get('subject', '默认学科')
        chapter = ai_result.get('chapter', '默认知识点')
        
        print(f"AI 解析结果：学科 = {subject}，知识点 = {chapter}")
        return subject, chapter
        
    except Exception as e:
        print(f"解析 AI 响应时出错：{e}。使用默认值。原始响应：\n{ai_response}")
        # 尝试从教学计划中提取前两行作为备选
        lines = teaching_plan.split('\n')
        subject = lines[0].strip() if len(lines) > 0 else "默认学科"
        chapter = lines[1].strip() if len(lines) > 1 else "默认知识点"
        return subject, chapter
def generate_PPT(subject=None, chapter=None, teaching_plan=None, teacher_name='AI', time=None, title=None, subtitle=None,
                 template=templatepath, ppt_filename=None, select='template'):  # 'ppts/template.pptx'
    if select == 'ai':
        pass
        # 连接AI PPT助手生成--尚未开发
    elif select == 'plan' and teaching_plan:
        if not subject or not chapter:
            # 调用解析函数
            subject, chapter = parse_teaching_plan(teaching_plan)
        
        content = get_template(subject, chapter, structure[filename], teaching_plan)
        if not isinstance(content, dict):
            raise ValueError('AI返回内容不可使用!')
        if title is None:
            title = subject
        if subtitle is None:
            subtitle = chapter[0] if chapter else '默认知识点'
        if time is None:
            time = datetime.now().strftime("%Y/%m/%d")

        prs = set_ppt_content(content, template, teacher_name, time, title, subtitle)

        if ppt_filename is None:
            unique_id = uuid.uuid4()  # 生成一个唯一的UUID
            ppt_filename = os.path.join(PPT_FOLDER, f"{subject}-{chapter[0]}教学PPT-{unique_id}.pptx")
        else:
            ppt_filename = os.path.join(PPT_FOLDER, ppt_filename if '.pptx' in ppt_filename else f'{ppt_filename}.pptx')

        # 保存PPT文件到本地
        prs.save(ppt_filename)
        print(f"PPT已保存为: {ppt_filename}")
    else:
        content = get_template(subject, chapter, structure[filename])
        if not isinstance(content, dict):
            raise ValueError('AI返回内容不可使用!')
        if title is None:
            title = subject
        if subtitle is None:
            subtitle = chapter[0] if chapter else '默认知识点'
        if time is None:
            time = datetime.now().strftime("%Y/%m/%d")

        prs = set_ppt_content(content, template, teacher_name, time, title, subtitle)

        if ppt_filename is None:
            unique_id = uuid.uuid4()  # 生成一个唯一的UUID
            ppt_filename = os.path.join(PPT_FOLDER, f"{subject}-{chapter[0]}教学PPT-{unique_id}.pptx")
        else:
            ppt_filename = os.path.join(PPT_FOLDER, ppt_filename if '.pptx' in ppt_filename else f'{ppt_filename}.pptx')

        # 保存PPT文件到本地
        prs.save(ppt_filename)
        print(f"PPT已保存为: {ppt_filename}")

    # 返回文件的存储路径
    return ppt_filename


def get_AI_content(template):
    client = OpenAI(api_key=key, base_url=api_url)
    response = client.chat.completions.create(
        model="deepseek-chat",  # 使用的模型
        messages=[
            {"role": "system", "content": "你是一个资深的教学设计专家，请用正式的语气回答用户的问题。"},
            {"role": "user", "content": template},
        ],
        stream=False  # 是否启用流式输出
    )
    # print(response)  # 如果需要调试，可以打印完整的响应
    content = response.choices[0].message.content  # 返回模型生成的回复内容

    return clean_content(content)


def clean_content(content):
    if 'json' in content:
        content = content.replace('json', '')
    while '`' in content:  # 或者if, replace完全匹配
        content = content.replace('`', '')
    if 'python' in content:
        content = content.replace('python', '')
    while '*' in content:
        content = content.replace('*', '')
    while '#' in content:
        content = content.replace('#', '')

    return content


def get_template(subject, chapter, structure, teaching_plan=None):
    template1 = lambda topic, number, type: (
        f"请围绕“{subject} {chapter} {topic}”设计 {number} 个教学{type}，"
        f"并以列表的形式返回各{type}名称，不需要序号，字数控制在10字以内，例如："
        f"[str1, str2, ...]") if type != '知识点' or teaching_plan is None else (
        f"请围绕“{subject} {chapter} {topic}”设计 {number} 个教学环节，必要时请包含一个互动环节，"
        f"并以列表的形式返回各环节名称，不需要序号，字数控制在10字以内，例如："
        f"[str1, str2, ...]"
        )
    # (f'请根据{theme}的内容设计{number}个教学{type}，并以Python列表的形式返回{type}名称')
    template2 = lambda theme, type: (f'请总结{subject} {chapter} {theme}，提取它的相关知识点，'
                                     f'并返回它的{type}，字数控制在30字左右，不可太多或太少；不需要返回标题和字数，仅返回内容')
    content = {}
    parts = []

    num1 = len(structure)
    if teaching_plan:
        parts = get_AI_content(f"请围绕以下教案内容，按序提取 {num1} 个关键教学环节：{teaching_plan}，"
                               f"并以列表的形式返回各环节名称，不需要序号，例如："
                               f"[str1, str2, ...]")
        print(parts)
    else:
        parts = get_AI_content(template1(f'', num1, '环节'))
        print(parts)
    
    text = parts.replace('，', ',').replace("'", '').replace('"', '')
    text = text.strip('[]').strip()
    parts = [item.strip() for item in text.split(',')]
    for i in range(num1):
        content[parts[i]] = {}

    for i in range(num1):
        # 生成概述
        summary = get_AI_content(template2(f'{parts[i]}的内容', '概述'))
        content[parts[i]]['概述'] = summary

        # 生成小标题
        part = structure[i]
        num2 = len(part)
        titles = get_AI_content(template1(parts[i], num2, '内容'))
        print(titles)
        
        text = titles.replace('，', ',').replace("'", '').replace('"', '')
        text = text.strip('[]').strip()
        titles = [item.strip() for item in text.split(',')]

        for j in range(num2):
            content[parts[i]][titles[j]] = {}

        for j in range(num2):
            subtitles = get_AI_content(template1(f'{parts[i]} {titles[j]}', part[j], '知识点'))
            print(subtitles)
            
            text = subtitles.replace('，', ',').replace("'", '').replace('"', '')
            text = text.strip('[]').strip()
            subtitles = [item.strip() for item in text.split(',')]
            for subtitle in subtitles:
                sub_content = get_AI_content(template2(f'{parts[i]} {titles[j]} {subtitle}的内容' if teaching_plan is None else
                                                       f'{parts[i]} {titles[j]}的{subtitle}环节', '教学内容'))  # 是否需要 {titles[j]}
                print(sub_content)
                content[parts[i]][titles[j]][subtitle] = sub_content

    return content