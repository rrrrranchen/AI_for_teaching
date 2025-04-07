# from sparkai.llm.llm import ChatSparkLLM
# from sparkai.core.messages import ChatMessage


# SPARKAI_URL = 'wss://spark-api.xf-yun.com/v4.0/chat'

# SPARKAI_APP_ID = 'fdc65b4b'
# SPARKAI_API_SECRET = 'ZDFmMzg5Yjk2ZTVhMmNhNjM2Nzc2ZmMz'
# SPARKAI_API_KEY = 'bd45dc727ac09175b09d3c2f5b2b993c'

# SPARKAI_DOMAIN = '4.0Ultra'

# def generate_lesson_plan(title, description):
#     return "教案示例"
    

#实现调用ai接口获取相关题目功能(参数)

import uuid
from openai import OpenAI
import time
import os


# 定义自己的 API Key
key = 'sk-b7550aa67ed840ffacb5ca051733802c'
api_url = "https://api.deepseek.com"  # DeepSeek 的 API 地址
# OpenAI 参数设置：API Key + API Interface (这里访问接口为 DeepSeek 的 API 地址)


# 逐字打印效果
def printChar(text, delay=0.1):
    for char in text:
        print(char, end='', flush=True)  # 使用 end='' 防止自动换行，flush=True 确保立即打印
        time.sleep(delay)
    print()  # 最后打印一个换行符


# 发送请求到 DeepSeek
def sendToDeepSeek(say):
    print('正在验证身份，请稍等....')
    # 请求接口并验证身份，创建客户端对象
    client = OpenAI(api_key=key, base_url=api_url)
    print('正在思考，请耐心等待...')
    # 发送请求数据并等待获取响应数据
    response = client.chat.completions.create(
        model="deepseek-chat",  # 使用的模型
        messages=[
            {"role": "system", "content": "你是一个专业的客服助手，请用正式的语气回答用户的问题。"},
            # {"role": "system", "content": "你是风趣幽默的客服，请用轻松幽默的语气回答用户的问题。"},
            {"role": "user", "content": say},
        ],
        stream=False  # 是否启用流式输出
    )
    # print(response)  # 如果需要调试，可以打印完整的响应
    return response.choices[0].message.content  # 返回模型生成的回复内容


# DeepSeek 问答环节
def callDeepSeek():
    # 主循环
    while True:
        myin = input('您请说：')  # 获取用户输入
        if myin == 'bye':  # 如果用户输入 "bye"，退出循环
            print('欢迎下次使用！再见！')
            break
        resp = sendToDeepSeek(myin)  # 发送用户输入到 DeepSeek 并获取回复
        printChar(resp)  # 逐字打印回复内容
        # print(resp)  # 如果需要直接打印完整回复，可以使用这行代码
        print('-----------------------------------------------------------')


# 测试接口
# callDeepSeek()


from pptx import Presentation, util
from datetime import datetime
from openai import OpenAI
from pptx import Presentation, util
from pptx.dml.color import RGBColor
import ast


def set_text_font(shape, font_name=None, font_size=None, font_color=None, font_bold=None, font_italic=None):
    text_frame = shape.text_frame

    # 遍历所有段落和文本运行
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            # 设置字体名称
            if font_name is not None:
                run.font.name = font_name
            # 设置字体大小（单位：磅）
            if font_size is not None:
                run.font.size = util.Pt(font_size)
            # 设置颜色
            if font_color is not None:  # NoneColor Type
                run.font.color.rgb = font_color
            # 设置加黑
            if font_bold is not None:  # False强制取消粗体, None继承父版样式
                run.font.bold = font_bold
            # 设置斜体
            if font_italic is not None:
                run.font.italic = font_italic
            # -> False显式取消加粗、斜体 -> if False... False & None 区别


def keep_text_font(shape, new_text):
    text_frame = shape.text_frame

    for paragraph in text_frame.paragraphs:
        # 创建新 Run 并继承格式（取第一个 Run 的格式）
        if paragraph.runs:
            # 清空所有 Run 的文本
            for run in paragraph.runs:
                run.text = ""

            first_run = paragraph.runs[0]
            new_run = paragraph.add_run()
            new_run.text = new_text

            # 复制字体属性
            new_run.font.name = first_run.font.name
            new_run.font.size = first_run.font.size
            new_run.font.bold = first_run.font.bold
            new_run.font.italic = first_run.font.italic
            # new_run.font.color.rgb = first_run.font.color.rgb
            # 检查颜色是否有效
            # if first_run.font.color is not None:  # and first_run.font.color.rgb -> 存在NoneType No RGB
            if (
                first_run.font.color is not None
                and hasattr(first_run.font.color, 'rgb')  # 关键检查
            ):
                new_run.font.color.rgb = first_run.font.color.rgb
            else:
                # 若颜色未设置，使用默认黑色或跳过
                new_run.font.color.rgb = RGBColor(0, 0, 0)  # 显式设置为黑色


def set_ppt_content(content, template, name, time, title, subtitle=None, keep_font=False,
                    font_name=None, font_size=None, font_bold=None, font_italic=None, font_color=None):
    prs = Presentation(template)

    # PPT模板的shapes全为文本框，这里可以省去判断步骤
    for slide in prs.slides:
        for shape in slide.shapes:
            if not (shape.has_text_frame and shape.text):
                continue
            text = shape.text
            # 设置标题 -> 防止副标题、小标题字样干扰
            if text == '标题':
                keep_text_font(shape, title)
                # if not keep_font: -> 设置字体样式... 这里保持为原格式
                set_text_font(shape, font_name='标准粗黑')
            # 设置副标题
            if text == '副标题':
                keep_text_font(shape, subtitle)
                set_text_font(shape, font_name='标准粗黑')
            # 设置老师姓名
            if '姓名' in text:
                keep_text_font(shape, f'主讲人：{name}')
                set_text_font(shape, font_name='等线')
            # 设置汇报日期
            if '日期' in text or '年月日' in text:
                keep_text_font(shape, f'时间：{time}')
                set_text_font(shape, font_name='等线')
            # 设置内容页
            if '小标题' in text:
                # shape.text = get_content(shape, content)
                # keep_text_font(shape, get_content(shape, content))
                sub_content, if_content = get_content(shape, content)
                keep_text_font(shape, sub_content)
                if not if_content:
                    set_text_font(shape, font_name='等线', font_bold=True)
            if not keep_font:
                set_text_font(shape, font_name, font_size, font_color, font_bold, font_italic)
    return prs


def get_content(shape, content):
    text = shape.text
    text = text.replace('小标题', '')

    if '概述' in text:
        i = int(text.replace('概述', '')) - 1
        return content[list(content.keys())[i]]['概述'], None  # False

    if_content = False

    if '内容' in text:
        text = text.replace('内容', '')
        if_content = True

    num = int(text.replace('.', ''))
    i, j, k = int(num / 100) - 1, int(num / 10) % 10 - 1, num % 10 - 1
    if j < 0:
        # return content.keys[k]
        key_lst = [k for k in content.keys()]
        return key_lst[k], if_content
    elif i < 0:
        # return content[content.keys[j]].keys[k]
        key_lst1 = list(content.keys())
        key_lst2 = [k for k in content[key_lst1[j]].keys() if k != '概述']  # list(content[key_lst1[j]].keys())
        return key_lst2[k], if_content
    # else:
    key_lst1 = list(content.keys())
    key_lst2 = [k for k in content[key_lst1[i]].keys() if k != '概述']
    key_lst3 = list(content[key_lst1[i]][key_lst2[j]].keys())
    # print(num, i+1, j+1, k+1, key_lst3)
    if if_content:
        return content[key_lst1[i]][key_lst2[j]][key_lst3[k]], if_content
    return key_lst3[k], if_content  # 输出if_content->调节小标题字体
    #     return content[content[content.keys[i]].keys[j]].values[k]
    # return content[content[content.keys[i]].keys[j]].keys[k]


project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPTTEMPLATE_FOLDER = os.path.join(project_root, 'static', 'template')
PPT_FOLDER=os.path.join(project_root,'static','uploads','ppts')
filename='template.pptx'
templatepath=os.path.join(PPTTEMPLATE_FOLDER, filename)
structure = {
'template.pptx': [
    [3, 3],
    [3, 3, 3],
    [3, 4, 3],
    [3, 3, 3, 3],
    [3, 3, 3]
    ]
}
def generate_PPT(subject, chapter, teaching_plan=None, teacher_name='AI', time=None, title=None, subtitle=None,
                 template=templatepath, ppt_filename=None, select='template'):  # 'ppts/template.pptx'
    if select == 'ai':
        pass
        # 连接AI PPT助手生成--尚未开发
    elif select == 'plan' and teaching_plan:
        content = get_template(subject, chapter, structure[filename], teaching_plan)
        if not isinstance(content, dict):
            raise ValueError('AI返回内容不可使用!')
        if title is None:
            title = subject
        if subtitle is None:
            subtitle = chapter
        if time is None:
            time = datetime.now().strftime("%Y/%m/%d")

        prs = set_ppt_content(content, template, teacher_name, time, title, subtitle)

        if ppt_filename is None:
            unique_id = uuid.uuid4()  # 生成一个唯一的UUID
            ppt_filename = os.path.join(PPT_FOLDER, f"{subject}-{chapter}教学PPT-{unique_id}.pptx")
        else:
            ppt_filename = os.path.join(PPT_FOLDER, ppt_filename if '.pptx' in ppt_filename else f'{ppt_filename}.pptx')

        # 保存PPT文件到本地
        prs.save(ppt_filename)
        print(f"PPT已保存为: {ppt_filename}")
    else:
        content = get_template(subject, chapter, structure[filename])
        if not isinstance(content, dict):
            raise ValueError('AI返回内容不可使用!')
        if title is None:
            title = subject
        if subtitle is None:
            subtitle = chapter
        if time is None:
            time = datetime.now().strftime("%Y/%m/%d")

        prs = set_ppt_content(content, template, teacher_name, time, title, subtitle)

        if ppt_filename is None:
            unique_id = uuid.uuid4()  # 生成一个唯一的UUID
            ppt_filename = os.path.join(PPT_FOLDER, f"{subject}-{chapter}教学PPT-{unique_id}.pptx")
        else:
            ppt_filename = os.path.join(PPT_FOLDER, ppt_filename if '.pptx' in ppt_filename else f'{ppt_filename}.pptx')

        # 保存PPT文件到本地
        prs.save(ppt_filename)
        print(f"PPT已保存为: {ppt_filename}")

    # 返回文件的存储路径
    return ppt_filename



def get_AI_content(template):
    client = OpenAI(api_key=key, base_url=api_url)
    response = client.chat.completions.create(
        model="deepseek-chat",  # 使用的模型
        messages=[
            {"role": "system", "content": "你是一个资深的教学设计专家，请用正式的语气回答用户的问题。"},
            {"role": "user", "content": template},
        ],
        stream=False  # 是否启用流式输出
    )
    # print(response)  # 如果需要调试，可以打印完整的响应
    content = response.choices[0].message.content  # 返回模型生成的回复内容

    return clean_content(content)


def clean_content(content):
    if 'json' in content:
        content = content.replace('json', '')
    while '`' in content:  # 或者if, replace完全匹配
        content = content.replace('`', '')
    if 'python' in content:
        content = content.replace('python', '')
    while '*' in content:
        content = content.replace('*', '')
    while '#' in content:
        content = content.replace('#', '')

    return content


def get_template(subject, chapter, structure, teaching_plan=None):
    template1 = lambda topic, number, type: (
        f"请围绕“{subject} {chapter} {topic}”设计 {number} 个教学{type}，"
        f"并以列表的形式返回各{type}名称，不需要序号，字数控制在10字以内，例如："
        f"[str1, str2, ...]") if type != '知识点' or teaching_plan is None else (
        f"请围绕“{subject} {chapter} {topic}”设计 {number} 个教学环节，必要时请包含一个互动环节，"
        f"并以列表的形式返回各环节名称，不需要序号，字数控制在10字以内，例如："
        f"[str1, str2, ...]"
        )
    # (f'请根据{theme}的内容设计{number}个教学{type}，并以Python列表的形式返回{type}名称')
    template2 = lambda theme, type: (f'请总结{subject} {chapter} {theme}，提取它的相关知识点，'
                                     f'并返回它的{type}，字数控制在30字左右，不可太多或太少；不需要返回标题和字数，仅返回内容')
    content = {}
    parts = []

    num1 = len(structure)
    if teaching_plan:
        parts = get_AI_content(f"请围绕以下教案内容，按序提取 {num1} 个关键教学环节：{teaching_plan}，"
                               f"并以列表的形式返回各环节名称，不需要序号，例如："
                               f"[str1, str2, ...]")
        print(parts)
    else:
        parts = get_AI_content(template1(f'', num1, '环节'))
        print(parts)
    
    text = parts.replace('，', ',').replace("'", '').replace('"', '')
    text = text.strip('[]').strip()
    parts = [item.strip() for item in text.split(',')]
    for i in range(num1):
        content[parts[i]] = {}

    for i in range(num1):
        # 生成概述
        summary = get_AI_content(template2(f'{parts[i]}的内容', '概述'))
        content[parts[i]]['概述'] = summary

        # 生成小标题
        part = structure[i]
        num2 = len(part)
        titles = get_AI_content(template1(parts[i], num2, '内容'))
        print(titles)
        
        text = titles.replace('，', ',').replace("'", '').replace('"', '')
        text = text.strip('[]').strip()
        titles = [item.strip() for item in text.split(',')]

        for j in range(num2):
            content[parts[i]][titles[j]] = {}

        for j in range(num2):
            subtitles = get_AI_content(template1(f'{parts[i]} {titles[j]}', part[j], '知识点'))
            print(subtitles)
            
            text = subtitles.replace('，', ',').replace("'", '').replace('"', '')
            text = text.strip('[]').strip()
            subtitles = [item.strip() for item in text.split(',')]
            for subtitle in subtitles:
                sub_content = get_AI_content(template2(f'{parts[i]} {titles[j]} {subtitle}的内容' if teaching_plan is None else
                                                       f'{parts[i]} {titles[j]}的{subtitle}环节', '教学内容'))  # 是否需要 {titles[j]}
                print(sub_content)
                content[parts[i]][titles[j]][subtitle] = sub_content

    return content


# print('开始：', datetime.now())
# generate_PPT('计算机网络', 'TCP-IP协议', teacher_name='丁力宏')
# plan = "# 计算机网络基础教案\n\n## 教学目标\n1. **知识目标**：学生能够准确描述计算机网络的三大功能（资源共享、信息传输、分布式处理），并能区分局域网、城域网和广域网的特点。\n2. **理解目标**：学生能够分析比较星型、总线型、环型三种网络拓扑结构的优缺点，准确率不低于90%。\n3. **应用目标**：学生能够根据给定的场景（如小型办公室、校园网等）选择合适的网络拓扑结构，并说明理由。\n4. **分析目标**：学生能够解释TCP/IP协议栈各层的功能及其相互关系，准确识别常见协议所属的层次。\n5. **评价目标**：学生能够评估不同网络设计方案的优势与局限，提出改进建议。\n\n## 教学重难点\n**重点**：\n1. 网络分类标准及各类网络的特点比较\n2. 三种主要拓扑结构的性能比较\n\n**难点**：\n1. 抽象的网络协议概念理解\n2. TCP/IP协议栈的分层逻辑\n\n**突破策略**：\n1. 使用类比教学法（如将网络协议比作语言交流规则）\n2. 通过网络模拟软件可视化数据流动过程\n3. 设计分层拼图游戏强化协议栈记忆\n\n## 教学内容\n1. **计算机网络基础概念**\n   - 定义：互联的自治计算机系统集合\n   - 三大功能模块：\n     * 资源共享（硬件、软件、数据）\n     * 信息传输（电子邮件、文件传输）\n     * 分布式处理（云计算、并行计算）\n\n2. **网络分类体系**\n   - 按覆盖范围：\n     * 局域网（LAN）：<1km，高带宽\n     * 城域网（MAN）：1-100km\n     * 广域网（WAN）：>100km，低带宽\n   - 按传输介质：有线/无线网络\n   - 按使用性质：公用/专用网络\n\n3. **拓扑结构专题**\n   - 星型结构：\n     * 中心节点压力大\n     * 单点故障影响\n   - 总线型结构：\n     * 冲突检测机制\n     * 扩展性限制\n   - 环型结构：\n     * 令牌传递机制\n     * 双向环冗余设计\n\n4. **TCP/IP协议栈**\n   - 四层模型：\n     * 应用层（HTTP/FTP/SMTP）\n     * 传输层（TCP/UDP）\n     * 网络层（IP/ICMP）\n     * 网络接口层\n\n## 教学时间安排\n| 教学环节 | 时间分配 |\n|---------|---------|\n| 导入活动 | 5分钟 |\n| 概念讲解 | 12分钟 |\n| 拓扑结构分析 | 10分钟 |\n| 协议栈探究 | 8分钟 |\n| 综合应用 | 7分钟 |\n| 总结提升 | 3分钟 |\n\n## 教学过程\n\n### 1. 导入环节（5分钟）\n**方法**：情境设问+实物演示  \n**活动**：  \n- 展示校园网拓扑图（互动环节1：快速问答）  \n  * 教师：\"当你在图书馆下载教学楼服务器上的文件时，数据经过哪些设备？\"  \n  * 学生抢答，教师用激光笔指示图中路径  \n**工具**：电子白板展示3D网络示意图  \n**预期**：激发兴趣，建立具象认知  \n\n### 2. 概念讲授（12分钟）\n**方法**：对比表格+案例教学  \n**活动**：  \n- 分发三种网络类型的对比表格（覆盖范围/典型应用/传输速率）  \n- 播放银行系统广域网案例视频  \n**学生行为**：小组填写表格关键数据  \n**教师行为**：  \n1. 用不同颜色标注表格重点  \n2. 演示Wireshark抓包过程（互动环节2：实时观测）  \n**材料**：预录制的数据包分析视频  \n\n### 3. 拓扑结构分析（10分钟）\n**方法**：模拟实验+角色扮演  \n**活动**：（互动环节3：拓扑模拟游戏）  \n- 将教室划分为三个区域，分别代表不同拓扑  \n- 学生分组扮演数据包，按规则\"传输\"笔记本  \n- 记录各结构传输耗时和故障影响范围  \n**工具**：计时器、故障情景卡片（如\"中心交换机断电\"）  \n**预期成果**：生成各拓扑性能对比报告  \n\n### 4. 协议栈探究（8分钟）\n**方法**：拼图竞赛+动画解析  \n**活动**：  \n- 分组拼装协议栈磁贴（应用层→物理层）  \n- 观看TCP三次握手动画（慢速播放+分段讲解）  \n**重点强调**：  \n- 协议封装/解封装过程  \n- 端口号与IP地址的协同作用  \n\n### 5. 综合应用（7分钟）\n**案例**：设计咖啡馆无线网络  \n**要求**：  \n1. 选择拓扑结构并说明理由  \n2. 列出需要的协议（至少3层）  \n3. 预测可能出现的网络问题  \n**展示方式**：小组代表使用思维导图软件汇报  \n\n### 6. 小结提升（3分钟）\n**方法**：概念地图填空  \n**活动**：  \n- 投影不完整的知识框架图  \n- 学生接力补充关键术语（网关、MTU、CSMA/CD等）  \n**升华提问**：\"为什么OSI七层模型在实际中较少使用？\"  \n\n## 课后作业\n**基础题**：  \n1. 绘制双绞线连接的星型拓扑图，标注关键设备  \n2. 匹配协议与层次（HTTP→应用层）  \n\n**拓展题**：  \n1. 分析智能家居系统适合的网络类型，需考虑：  \n   - 设备数量（20+ IoT设备）  \n   - 数据传输特点（小数据包、高频率）  \n2. 用Wireshark捕获一次网页访问过程，识别至少3种协议  \n\n**挑战题**（选做）：  \n设计校园宿舍区网络改造方案，要求：  \n- 比较有线与无线方案的优劣  \n- 估算500个终端同时在线所需带宽  \n- 列出可能采用的QoS策略"
# generate_PPT('计算机网络', '', teacher_name='朱钟', teaching_plan=plan, select='plan')
# print('结束：', datetime.now())

