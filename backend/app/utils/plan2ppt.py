from openai import OpenAI
from app.config import Config
import time
import os

# 定义自己的 API Key
key = Config.DEEPSEEK_API_KEY
api_url = "https://api.deepseek.com"  # DeepSeek 的 API 地址


from pptx import Presentation, util
from datetime import datetime
from openai import OpenAI
from pptx import Presentation, util
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import ast


def set_text_font(shape, font_name=None, font_size=None, font_color=None, font_bold=None, font_italic=None):
    text_frame = shape.text_frame

    # 遍历所有段落和文本运行
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            # 设置字体名称
            if font_name is not None:
                run.font.name = font_name
            # 设置字体大小（单位：磅）
            if font_size is not None:
                run.font.size = util.Pt(font_size)
            # 设置颜色
            if font_color is not None:  # NoneColor Type
                run.font.color.rgb = font_color
            # 设置加黑
            if font_bold is not None:  # False强制取消粗体, None继承父版样式
                run.font.bold = font_bold
            # 设置斜体
            if font_italic is not None:
                run.font.italic = font_italic
            # -> False显式取消加粗、斜体 -> if False... False & None 区别


def keep_text_font(shape, new_text):
    text_frame = shape.text_frame

    for paragraph in text_frame.paragraphs:
        # 创建新 Run 并继承格式（取第一个 Run 的格式）
        if paragraph.runs:
            # 清空所有 Run 的文本
            for run in paragraph.runs:
                run.text = ""

            first_run = paragraph.runs[0]
            new_run = paragraph.add_run()
            new_run.text = new_text

            # 复制字体属性
            new_run.font.name = first_run.font.name
            new_run.font.size = first_run.font.size
            new_run.font.bold = first_run.font.bold
            new_run.font.italic = first_run.font.italic
            # new_run.font.color.rgb = first_run.font.color.rgb
            # 检查颜色是否有效
            # if first_run.font.color is not None:  # and first_run.font.color.rgb -> 存在NoneType No RGB
            if (
                    first_run.font.color is not None
                    and hasattr(first_run.font.color, 'rgb')  # 关键检查
            ):
                new_run.font.color.rgb = first_run.font.color.rgb
            else:
                # 若颜色未设置，使用默认黑色或跳过
                new_run.font.color.rgb = RGBColor(0, 0, 0)  # 显式设置为黑色


def list_format(content):
    text = (content.replace('json', '').replace('python', '')
            .replace('`', '').replace('*', '').replace('#', '')
            .replace('，', ',').replace("'", '').replace('"', ''))
    text = text.strip('[]').strip()
    content = [item.strip() for item in text.split(',')]

    return content


def fromAI(say):
    if isinstance(say, str):
        client = OpenAI(api_key=key, base_url=api_url)
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": "你是一个资深的教学设计专家，请用正式的语气简洁地回答用户的问题。"},
                {"role": "user", "content": say}
            ],
            stream=False
        )
        return response.choices[0].message.content

    client = OpenAI(api_key=key, base_url=api_url)
    response = client.chat.completions.create(
        model="deepseek-chat",
        messages=say,
        stream=False
    )
    return response.choices[0].message.content


def group_shape(shape, title, subtitle, name, time, content):
    # 检查是否是嵌套文本框
    for sub_shape in shape.shapes:
        if sub_shape.shape_type == 6:
            group_shape(sub_shape, title, subtitle, name, time, content)
            continue
        # 检查是否是文本框
        if sub_shape.has_text_frame:
            text_frame = sub_shape.text_frame
            if text_frame.text:
                fill_text(sub_shape, text_frame.text, title, subtitle, name, time, content)
        # 检查是否是图形
        elif sub_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if hasattr(sub_shape, "text_frame"):
                text_frame = sub_shape.text_frame
                if text_frame.text:
                    fill_text(sub_shape, text_frame.text, title, subtitle, name, time, content)


def fill_text(shape, text, title, subtitle, name, time, content):
    text = text.replace(' ', '')
    # 设置标题 -> 防止副标题、小标题字样干扰
    if text == '标题':
        # shape.text = title
        keep_text_font(shape, title)
        set_text_font(shape, font_name='标准粗黑')
    # 设置副标题
    if text == '副标题':
        # shape.text = subtitle
        keep_text_font(shape, subtitle)
        set_text_font(shape, font_name='标准粗黑')
    # 设置老师姓名
    if '姓名' in text:
        # shape.text = f'主讲人：{name}'
        keep_text_font(shape, f'主讲人：{name}')
        set_text_font(shape, font_name='等线')
    # 设置汇报日期
    if '日期' in text or '年月日' in text:
        # shape.text = f'时间：{time}'
        keep_text_font(shape, f'时间：{time}')
        set_text_font(shape, font_name='等线')
    # 设置内容页
    if '小标题' in text:
        sub_content, if_content = find_content(shape, content)
        # shape.text = sub_content
        keep_text_font(shape, sub_content)
        if if_content:
            set_text_font(shape, font_name='等线')
        else:
            set_text_font(shape, font_name='等线', font_bold=True)


def fill_ppt(subject, content, template, name, subtitle):
    prs = Presentation(template)
    title = subject
    if subtitle is None:
        subtitle = '教学示例'
    time = datetime.now().strftime("%Y/%m/%d")

    for slide in prs.slides:
        for shape in slide.shapes:
            # if not (shape.has_text_frame and shape.text):
            #     continue
            # text = shape.text
            # 检查是否是文本框
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if text_frame.text:
                    fill_text(shape, text_frame.text, title, subtitle, name, time, content)
            # 检查是否是图形
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if hasattr(shape, "text_frame"):
                    text_frame = shape.text_frame
                    if text_frame.text:
                        fill_text(shape, text_frame.text, title, subtitle, name, time, content)
            elif shape.shape_type == 6:  # 6 代表 GroupShape
                group_shape(shape, title, subtitle, name, time, content)
            
    return prs


def find_content(shape, content):
    text = shape.text
    text = text.replace('小标题', '')

    if '概述' in text:
        i = int(text.replace('概述', '')) - 1
        return content[list(content.keys())[i]]['概述'], False

    if_content = False

    if '内容' in text:
        text = text.replace('内容', '')
        if_content = True

    num = int(text.replace('.', ''))
    i, j, k = int(num / 100) - 1, int(num / 10) % 10 - 1, num % 10 - 1
    if j < 0:
        # return content.keys[k]
        key_lst = [k for k in content.keys()]
        return key_lst[k], if_content
    elif i < 0:
        # return content[content.keys[j]].keys[k]
        key_lst1 = list(content.keys())
        key_lst2 = [k for k in content[key_lst1[j]].keys() if k != '概述']  # list(content[key_lst1[j]].keys())
        return key_lst2[k], if_content
    # else:
    key_lst1 = list(content.keys())
    key_lst2 = [k for k in content[key_lst1[i]].keys() if k != '概述']
    key_lst3 = list(content[key_lst1[i]][key_lst2[j]].keys())
    # print(num, i+1, j+1, k+1, key_lst3)
    if if_content:
        return content[key_lst1[i]][key_lst2[j]][key_lst3[k]], if_content
    return key_lst3[k], if_content  # 输出if_content->调节小标题字体
    #     return content[content[content.keys[i]].keys[j]].values[k]
    # return content[content[content.keys[i]].keys[j]].keys[k]


def get_content(plan, structure):
    say0 = lambda x: (f"请根据以下教案:"
                      f"{x},"
                      f"仅返回它的学科名")
    say1 = lambda x, n: (f"请根据以下教案:"
                         f"{x},"
                         f"总结提取{n}个教学环节，按以下格式输出："
                         f"['第一个环节', '第二个环节', ...],"
                         f"要求：不要备注，每个环节名称不超过10个字。")  # 、不要标题
    say2 = lambda x, t, n, m: (f"请结合{x}课程的实际教学，"
                               f"请围绕'{t}'总结提取返回{n}个不同的教学内容，按以下格式输出："
                               f"['第一个内容', '第二个内容', ...],"
                               f"要求：不要备注、不要标题，每个内容名称不超过10个字。"
                               f"{m}")
    say3 = lambda x, t, tt, w: (f"请结合{x}课程的实际教学，"
                                f"请围绕'{t}'主题 返回它的{tt}，"
                                f"要求：不要备注、不要标题，仅返回内容，需要{w}个字左右。")
    content = {}
    num1 = len(structure)
    subject = fromAI(say0(plan))
    parts = list_format(fromAI(say1(plan, num1)))
    print(subject)
    print(parts)

    for i in range(num1):
        content[parts[i]] = {}

    for i in range(num1):
        # 生成概述
        summary = fromAI(say3(subject, parts[i], '概述', 35))
        content[parts[i]]['概述'] = summary
        print(summary)

        # 生成小标题
        part = structure[i]
        num2 = len(part)
        titles = list_format(fromAI(say2(subject, parts[i], num2, '')))
        print(titles)

        for j in range(num2):
            content[parts[i]][titles[j]] = {}

        for j in range(num2):
            m = '必要时可包含一个互动内容' if j == num2 - 1 else ''
            subtitles = list_format(fromAI(say2(subject, f'{parts[i]}的{titles[j]}阶段', part[j], m)))
            print(subtitles)

            for subtitle in subtitles:
                sub_content = fromAI(say3(subject, f'{subtitle}阶段', '教学内容', 30))
                print(sub_content)
                content[parts[i]][titles[j]][subtitle] = sub_content

    return subject, content


project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPTTEMPLATE_FOLDER = os.path.join(project_root, 'static', 'template')
PPT_FOLDER = os.path.join(project_root, 'static', 'uploads', 'ppts')
filename = 'template7.pptx'
templatepath = os.path.join(PPTTEMPLATE_FOLDER, filename)
structure = {
    'template1.pptx': [
        [3, 3],
        [3, 3, 3],
        [3, 4, 3],
        [3, 3, 3, 3],
        [3, 3, 3]
    ],
    'template2.pptx': [
        [3],
        [3],
        [3, 3],
        [3, 3, 3]
    ],
    'template3.pptx': [
        [3, 4],
        [4],
        [4, 3],
        [4]
    ],
    'template4.pptx': [
        [5, 4, 4, 4],
        [4, 5, 3],
        [4, 4, 4, 6],
        [3, 4, 3, 3]
    ],
    'template5.pptx': [
        [3, 4, 4, 5],
        [4, 5, 4, 4],
        [3, 5, 4],
        [3, 6, 3]
    ],
    'template6.pptx': [
        [5, 4, 4, 4],
        [4, 3, 3],
        [3, 4],
        [4, 4, 4, 4]
    ],
    'template7.pptx': [
        [4, 6, 6, 6],
        [5, 4, 4, 4],
        [4, 4, 4, 4],
        [4, 3, 3, 3]
    ]
}


def generate_PPT(teaching_plan, template, teacher_name, subtitle=None):  # 设置副标题，大标题统一为 学科名
        subject, content = get_content(teaching_plan, structure[filename])
        print(subject)
        print(content)
        prs = fill_ppt(subject, content, template, teacher_name, subtitle)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        ppt_filename = os.path.join(PPT_FOLDER, f"{subject}教学PPT{timestamp}.pptx")
        
        # 保存PPT文件到本地
        prs.save(ppt_filename)
        print(f"PPT已保存为: {ppt_filename}")
        return ppt_filename


def plan2dict(plan, teacher, template=filename):
    templatepath = os.path.join(PPTTEMPLATE_FOLDER, filename)
    stt = structure[template]
    num1 = len(stt)

    say = lambda x, n: (f"请根据以下教案:"
                        f"{x},"
                        f"总结提取{n}个教学环节，按以下格式输出："
                        f"['第一个环节', '第二个环节', ...], 每个环节名称不超过6个字。")
    say1 = lambda x: (f"请根据以下教案:"
                      f"{x},"
                      f"提取和返回它的学科名")
    say2 = lambda t, n: (f"在前问的基础上，结合本学科的实际教学，"
                         f"请围绕'{t}'总结提取返回{n}个教学内容，按以下格式输出："
                         f"['第一个内容', '第二个内容', ...], 每个内容名称不超过7个字，必要时可包含一个互动内容。")
    say3 = lambda t, tt, w: (f"在前问的基础上，结合本学科的实际教学，"
                             f"请围绕'{t}'返回{tt}，不超过{w}个字。")

    parts = list_format(fromAI(say(plan, num1)))
    print(parts)
    says = [{"role": "system", "content": "你是一个专业的客服助手，请用正式的语气回答用户的问题，每次回复请按以下格式返回：{1: '第一题回复', 2: '第二题回复', ...}，每次回复不超过10字"}]

    dict1 = {}
    while len(dict1) != 2 * num1 + 1:
        says1 = says
        for i in range(num1):
            says1.append({"role": "user", "content": say1(plan)})
            says1.append({"role": "user", "content": say2(parts[i], len(stt[i]))})
            says1.append({"role": "user", "content": say3(parts[i], '它的概述', 30)})

        dict1 = eval(fromAI(says1))

    vals = list(dict1.values())
    subject = vals[0]
    titles, sums = [], []
    for v in vals:
        if v == subject:
            continue
        if isinstance(v, str):
            sums.append(v)
        elif isinstance(v, list):
            titles.append(v)


#plan = "# 计算机网络基础教案\n\n## 教学目标\n1. **知识目标**：学生能够准确描述计算机网络的三大功能（资源共享、信息传输、分布式处理），并能区分局域网、城域网和广域网的特点。\n2. **理解目标**：学生能够分析比较星型、总线型、环型三种网络拓扑结构的优缺点，准确率不低于90%。\n3. **应用目标**：学生能够根据给定的场景（如小型办公室、校园网等）选择合适的网络拓扑结构，并说明理由。\n4. **分析目标**：学生能够解释TCP/IP协议栈各层的功能及其相互关系，准确识别常见协议所属的层次。\n5. **评价目标**：学生能够评估不同网络设计方案的优势与局限，提出改进建议。\n\n## 教学重难点\n**重点**：\n1. 网络分类标准及各类网络的特点比较\n2. 三种主要拓扑结构的性能比较\n\n**难点**：\n1. 抽象的网络协议概念理解\n2. TCP/IP协议栈的分层逻辑\n\n**突破策略**：\n1. 使用类比教学法（如将网络协议比作语言交流规则）\n2. 通过网络模拟软件可视化数据流动过程\n3. 设计分层拼图游戏强化协议栈记忆\n\n## 教学内容\n1. **计算机网络基础概念**\n   - 定义：互联的自治计算机系统集合\n   - 三大功能模块：\n     * 资源共享（硬件、软件、数据）\n     * 信息传输（电子邮件、文件传输）\n     * 分布式处理（云计算、并行计算）\n\n2. **网络分类体系**\n   - 按覆盖范围：\n     * 局域网（LAN）：<1km，高带宽\n     * 城域网（MAN）：1-100km\n     * 广域网（WAN）：>100km，低带宽\n   - 按传输介质：有线/无线网络\n   - 按使用性质：公用/专用网络\n\n3. **拓扑结构专题**\n   - 星型结构：\n     * 中心节点压力大\n     * 单点故障影响\n   - 总线型结构：\n     * 冲突检测机制\n     * 扩展性限制\n   - 环型结构：\n     * 令牌传递机制\n     * 双向环冗余设计\n\n4. **TCP/IP协议栈**\n   - 四层模型：\n     * 应用层（HTTP/FTP/SMTP）\n     * 传输层（TCP/UDP）\n     * 网络层（IP/ICMP）\n     * 网络接口层\n\n## 教学时间安排\n| 教学环节 | 时间分配 |\n|---------|---------|\n| 导入活动 | 5分钟 |\n| 概念讲解 | 12分钟 |\n| 拓扑结构分析 | 10分钟 |\n| 协议栈探究 | 8分钟 |\n| 综合应用 | 7分钟 |\n| 总结提升 | 3分钟 |\n\n## 教学过程\n\n### 1. 导入环节（5分钟）\n**方法**：情境设问+实物演示  \n**活动**：  \n- 展示校园网拓扑图（互动环节1：快速问答）  \n  * 教师：\"当你在图书馆下载教学楼服务器上的文件时，数据经过哪些设备？\"  \n  * 学生抢答，教师用激光笔指示图中路径  \n**工具**：电子白板展示3D网络示意图  \n**预期**：激发兴趣，建立具象认知  \n\n### 2. 概念讲授（12分钟）\n**方法**：对比表格+案例教学  \n**活动**：  \n- 分发三种网络类型的对比表格（覆盖范围/典型应用/传输速率）  \n- 播放银行系统广域网案例视频  \n**学生行为**：小组填写表格关键数据  \n**教师行为**：  \n1. 用不同颜色标注表格重点  \n2. 演示Wireshark抓包过程（互动环节2：实时观测）  \n**材料**：预录制的数据包分析视频  \n\n### 3. 拓扑结构分析（10分钟）\n**方法**：模拟实验+角色扮演  \n**活动**：（互动环节3：拓扑模拟游戏）  \n- 将教室划分为三个区域，分别代表不同拓扑  \n- 学生分组扮演数据包，按规则\"传输\"笔记本  \n- 记录各结构传输耗时和故障影响范围  \n**工具**：计时器、故障情景卡片（如\"中心交换机断电\"）  \n**预期成果**：生成各拓扑性能对比报告  \n\n### 4. 协议栈探究（8分钟）\n**方法**：拼图竞赛+动画解析  \n**活动**：  \n- 分组拼装协议栈磁贴（应用层→物理层）  \n- 观看TCP三次握手动画（慢速播放+分段讲解）  \n**重点强调**：  \n- 协议封装/解封装过程  \n- 端口号与IP地址的协同作用  \n\n### 5. 综合应用（7分钟）\n**案例**：设计咖啡馆无线网络  \n**要求**：  \n1. 选择拓扑结构并说明理由  \n2. 列出需要的协议（至少3层）  \n3. 预测可能出现的网络问题  \n**展示方式**：小组代表使用思维导图软件汇报  \n\n### 6. 小结提升（3分钟）\n**方法**：概念地图填空  \n**活动**：  \n- 投影不完整的知识框架图  \n- 学生接力补充关键术语（网关、MTU、CSMA/CD等）  \n**升华提问**：\"为什么OSI七层模型在实际中较少使用？\"  \n\n## 课后作业\n**基础题**：  \n1. 绘制双绞线连接的星型拓扑图，标注关键设备  \n2. 匹配协议与层次（HTTP→应用层）  \n\n**拓展题**：  \n1. 分析智能家居系统适合的网络类型，需考虑：  \n   - 设备数量（20+ IoT设备）  \n   - 数据传输特点（小数据包、高频率）  \n2. 用Wireshark捕获一次网页访问过程，识别至少3种协议  \n\n**挑战题**（选做）：  \n设计校园宿舍区网络改造方案，要求：  \n- 比较有线与无线方案的优劣  \n- 估算500个终端同时在线所需带宽  \n- 列出可能采用的QoS策略"
#generate_PPT(plan, templatepath, '刘俊涛')
